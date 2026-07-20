"""Export a deck to a real .pptx — color-coded by pedagogical move, mirroring the
on-screen slides (accent bar, move badge, tinted background, colored title)."""

import io

import httpx
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .models import Deck, Slide
from .render import slide_lines, slide_title

_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_INK = RGBColor(0x1A, 0x1F, 0x2B)
_COVER = RGBColor(0x2E, 0x40, 0xC4)

# Per-move (bold, pale) — the same palette the UI uses.
_MOVE: dict[str, tuple[str, str]] = {
    "hook": ("F9800A", "FFF3E2"),
    "concept": ("3A86FF", "E8F1FF"),
    "check-for-understanding": ("2EA043", "E6F8EA"),
    "exit-ticket": ("8338EC", "F3EBFF"),
    "video": ("E5385B", "FFE9EF"),
    "sorting-game": ("0891B2", "CFFAFE"),
}
_CUSTOM = ("00897B", "E0F2F1")
_EMOJI = {
    "hook": "🎣", "concept": "💡", "check-for-understanding": "✅",
    "exit-ticket": "🎟️", "video": "🎬", "sorting-game": "🎮",
}


def _rgb(hex6: str) -> RGBColor:
    return RGBColor(int(hex6[0:2], 16), int(hex6[2:4], 16), int(hex6[4:6], 16))


def _fetch_image(url: str) -> io.BytesIO | None:
    try:
        r = httpx.get(url, timeout=6, follow_redirects=True)
        if r.status_code == 200 and r.content:
            return io.BytesIO(r.content)
    except Exception:
        return None
    return None


def _no_line(shape) -> None:
    shape.line.fill.background()
    shape.shadow.inherit = False


def _content_slide(prs: Presentation, slide: Slide) -> None:
    bold_hex, pale_hex = _MOVE.get(slide.type, _CUSTOM)
    bold, pale = _rgb(bold_hex), _rgb(pale_hex)
    W = prs.slide_width
    H = prs.slide_height
    s = prs.slides.add_slide(prs.slide_layouts[6])

    # Tinted background + a bold accent bar down the left edge.
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = pale
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = bold
    _no_line(bar)

    # Move badge (a colored pill).
    label = f"{_EMOJI.get(slide.type, '🧩')} {slide.type}"
    badge = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.45),
        Inches(min(5.0, 1.0 + 0.11 * len(label))), Inches(0.42),
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = bold
    _no_line(badge)
    bp = badge.text_frame.paragraphs[0]
    bp.text = label
    bp.font.size = Pt(12)
    bp.font.bold = True
    bp.font.color.rgb = _WHITE
    bp.alignment = PP_ALIGN.CENTER

    img = _fetch_image(str(slide.content.get("image_url", "")))
    body_w = Inches(7.2) if img else Inches(12.0)
    if img:
        try:
            s.shapes.add_picture(img, Inches(8.3), Inches(1.9), width=Inches(4.5))
        except Exception:
            pass

    # Title (in the move color).
    tbox = s.shapes.add_textbox(Inches(0.6), Inches(1.1), body_w, Inches(1.0))
    tp = tbox.text_frame.paragraphs[0]
    tbox.text_frame.word_wrap = True
    tp.text = slide_title(slide)
    tp.font.size = Pt(30)
    tp.font.bold = True
    tp.font.color.rgb = bold

    # Body.
    cbox = s.shapes.add_textbox(Inches(0.6), Inches(2.1), body_w, Inches(4.9))
    ctf = cbox.text_frame
    ctf.word_wrap = True
    first = True
    for line in slide_lines(slide):
        para = ctf.paragraphs[0] if first else ctf.add_paragraph()
        para.text = line
        para.font.size = Pt(16)
        para.font.color.rgb = _INK
        para.space_after = Pt(9)
        first = False


def build_pptx(deck: Deck) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Vibrant cover: full-bleed accent with white text.
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    cover.background.fill.solid()
    cover.background.fill.fore_color.rgb = _COVER
    box = cover.shapes.add_textbox(Inches(1.0), Inches(2.6), Inches(11.3), Inches(2.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = deck.subject.title()
    p.font.size = Pt(46)
    p.font.bold = True
    p.font.color.rgb = _WHITE
    sub = tf.add_paragraph()
    sub.text = f"{deck.grade}  ·  {len(deck.slides)} slides"
    sub.font.size = Pt(22)
    sub.font.color.rgb = _WHITE

    for slide in deck.slides:
        _content_slide(prs, slide)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
